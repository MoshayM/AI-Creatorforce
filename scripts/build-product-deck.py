# Builds docs/CreatorForce-Product-Overview.pptx — product + go-to-market deck.
# Grounded in the repo docs (project.md, roadmap.md, features.md,
# monetization-framework.md) and the current implementation state.
# Run: python scripts/build-product-deck.py

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────────────────
INK = RGBColor(0x12, 0x14, 0x1C)      # near-black background
PANEL = RGBColor(0x1C, 0x1F, 0x2B)    # card background
ACCENT = RGBColor(0xF6, 0x3B, 0x5C)   # coral red (creator energy)
ACCENT2 = RGBColor(0x4C, 0xC9, 0xF0)  # cyan (tech)
GOLD = RGBColor(0xFF, 0xC5, 0x4C)
WHITE = RGBColor(0xF5, 0xF6, 0xFA)
MUTED = RGBColor(0x9A, 0xA3, 0xB5)
GREEN = RGBColor(0x51, 0xC9, 0x8B)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        font="Segoe UI", anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return box


def bullets(s, x, y, w, h, items, size=15, color=WHITE, gap=6, lead_color=ACCENT):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, rest = item
            r = p.add_run(); r.text = "▸ "; r.font.color.rgb = lead_color; r.font.size = Pt(size); r.font.bold = True
            r = p.add_run(); r.text = lead; r.font.color.rgb = color; r.font.size = Pt(size); r.font.bold = True
            r = p.add_run(); r.text = " — " + rest; r.font.color.rgb = MUTED; r.font.size = Pt(size)
        else:
            r = p.add_run(); r.text = "▸ "; r.font.color.rgb = lead_color; r.font.size = Pt(size); r.font.bold = True
            r = p.add_run(); r.text = item; r.font.color.rgb = color; r.font.size = Pt(size)
        for r in p.runs:
            r.font.name = "Segoe UI"
    return box


def card(s, x, y, w, h, fill=PANEL, line_color=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = 0.06
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    if line_color:
        c.line.color.rgb = line_color
        c.line.width = Pt(1)
    else:
        c.line.fill.background()
    c.shadow.inherit = False
    return c


def chip(s, x, y, w, h, text, fill=ACCENT, tcolor=WHITE, size=13, bold=True):
    c = card(s, x, y, w, h, fill=fill)
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor; r.font.name = "Segoe UI"
    return c


def header(s, kicker, title):
    txt(s, Inches(0.6), Inches(0.32), Inches(9), Inches(0.4), kicker.upper(),
        size=13, color=ACCENT, bold=True)
    txt(s, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.9), title,
        size=30, color=WHITE, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.38), Inches(1.4), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False


def arrow(s, x, y, w=Inches(0.32)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background(); a.shadow.inherit = False
    return a


def footer(s, n):
    txt(s, Inches(0.6), Inches(7.08), Inches(6), Inches(0.3),
        "AI CreatorForce — Product & Market Overview", size=10, color=MUTED)
    txt(s, Inches(12.4), Inches(7.08), Inches(0.6), Inches(0.3), str(n), size=10,
        color=MUTED, align=PP_ALIGN.RIGHT)


# ═══ 1. Title ══════════════════════════════════════════════════════════════════
s = slide()
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(-2.4), Inches(8), Inches(8))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x25, 0x14, 0x1E); glow.line.fill.background(); glow.shadow.inherit = False
chip(s, Inches(0.7), Inches(1.75), Inches(2.6), Inches(0.42), "PRODUCT + GTM BRIEF", fill=PANEL, tcolor=ACCENT2)
txt(s, Inches(0.65), Inches(2.3), Inches(11.5), Inches(1.6),
    "AI CreatorForce", size=60, bold=True)
txt(s, Inches(0.7), Inches(3.45), Inches(11.5), Inches(0.7),
    "The YouTube Content Operating System", size=26, color=ACCENT)
txt(s, Inches(0.7), Inches(4.35), Inches(10.8), Inches(1.1),
    "AI does the production work — research, scripts, editing, shorts, SEO.\n"
    "Humans keep the judgment — every publish passes a compliance gate and a human approval.",
    size=16, color=MUTED, line_spacing=1.15)
txt(s, Inches(0.7), Inches(6.5), Inches(11), Inches(0.4),
    "July 2026  ·  Platform status: feature-complete core, production-hardening phase", size=13, color=MUTED)

# ═══ 2. The problem ════════════════════════════════════════════════════════════
s = slide(); header(s, "Why this exists", "Creators drown in production work, not ideas")
data = [
    ("10–20 hrs", "to produce one quality long-form video — research, script, edit, thumbnail, SEO", ACCENT),
    ("3+ tools", "stitched together per video: an SEO tool, a clip tool, an editor, a scheduler", ACCENT2),
    ("1 strike", "a monetization or copyright misstep can demonetize a channel overnight", GOLD),
    ("70%+", "of long-form value is never repurposed into Shorts, posts, or newsletters", GREEN),
]
x = Inches(0.6)
for big, small, col in data:
    c = card(s, x, Inches(1.9), Inches(2.95), Inches(2.6))
    txt(s, x + Inches(0.25), Inches(2.2), Inches(2.5), Inches(0.8), big, size=34, color=col, bold=True)
    txt(s, x + Inches(0.25), Inches(3.1), Inches(2.5), Inches(1.3), small, size=13, color=MUTED, line_spacing=1.1)
    x += Inches(3.15)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6),
    "The gap in the market: existing tools optimize one step (clips, SEO, editing). Nobody owns the "
    "whole workflow with monetization safety built in.\nCreatorForce is the workflow — idea to published, "
    "compliance-gated at the core, human-approved at the end.",
    size=16, color=WHITE, line_spacing=1.25)
footer(s, 2)

# ═══ 3. What the software is ═══════════════════════════════════════════════════
s = slide(); header(s, "Software overview", "One platform, four surfaces, one AI agent workforce")
surfaces = [
    ("Web App", "Next.js dashboard — projects, Shorts Studio, video editor, approvals, analytics, wallet", ACCENT),
    ("AI Agent Suite", "15+ specialized agents (research, script, fact-check, compliance, SEO, media) orchestrated by a supervisor over a job queue", ACCENT2),
    ("Automation", "Per-channel autopilot: auto-import, auto-analyze, paced auto-publish of pre-approved content", GOLD),
    ("Developer API", "Scoped API keys, webhooks, rate-limited public endpoints for programmatic pipelines", GREEN),
]
y = Inches(1.75)
for name, desc, col in surfaces:
    card(s, Inches(0.6), y, Inches(7.5), Inches(1.15))
    txt(s, Inches(0.85), y + Inches(0.12), Inches(2.4), Inches(0.5), name, size=17, color=col, bold=True)
    txt(s, Inches(3.15), y + Inches(0.12), Inches(4.8), Inches(0.95), desc, size=12.5, color=MUTED, line_spacing=1.05)
    y += Inches(1.3)
c = card(s, Inches(8.4), Inches(1.75), Inches(4.3), Inches(5.0), line_color=ACCENT2)
txt(s, Inches(8.65), Inches(1.95), Inches(3.8), Inches(0.4), "SCOPE", size=13, color=ACCENT2, bold=True)
bullets(s, Inches(8.65), Inches(2.35), Inches(3.9), Inches(4.2), [
    ("In", "idea → research → script → media → shorts → SEO → publish → analyze"),
    ("In", "teams/orgs, budgets, credits, referrals, marketplace"),
    ("Out", "spam or content-farm tooling"),
    ("Out", "publishing without human approval"),
    ("Out", "copyright / disclosure evasion"),
], size=13, gap=10)
txt(s, Inches(8.65), Inches(6.0), Inches(3.9), Inches(0.7),
    "Non-goals are enforced in code, not policy docs — no publish path bypasses the compliance agent.",
    size=11.5, color=MUTED, line_spacing=1.1)
footer(s, 3)

# ═══ 4. Core flow end-to-end (long-form) ═══════════════════════════════════════
s = slide(); header(s, "Core flow · end to end", "Long-form pipeline: idea in, upload-ready video out")
stages = [
    ("Research", "sources +\ntrends", ACCENT2), ("Script", "voice-matched\ndraft", ACCENT2),
    ("Fact Check", "claims traced\nto sources", ACCENT2), ("Compliance", "HARD GATE\nblocks pipeline", ACCENT),
    ("SEO / Meta", "title, tags,\ndescription", ACCENT2), ("Media Gen", "voice, images,\nmusic, video", ACCENT2),
    ("Render", "ffmpeg\ntimeline", ACCENT2), ("Approval", "HUMAN\nDECIDES", GOLD), ("Publish", "YouTube\nData API", GREEN),
]
x = Inches(0.42); y = Inches(2.1)
for i, (name, sub, col) in enumerate(stages):
    w = Inches(1.24)
    c = card(s, x, y, w, Inches(1.5), fill=PANEL, line_color=col)
    txt(s, x + Inches(0.04), y + Inches(0.12), Inches(1.16), Inches(0.5), name, size=12.5, color=col, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.04), y + Inches(0.62), Inches(1.16), Inches(0.8), sub, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < len(stages) - 1:
        arrow(s, x + w + Inches(0.015), y + Inches(0.6), w=Inches(0.16))
    x += w + Inches(0.19)
bullets(s, Inches(0.6), Inches(4.1), Inches(12.0), Inches(2.6), [
    ("Async by design", "every stage is a BullMQ job with live progress over websockets; stages self-skip when output already exists, so re-runs only fill gaps"),
    ("Two independent gates", "the AI compliance audit blocks the pipeline; a human approval blocks the upload. Neither can be bypassed by any code path"),
    ("Validated at every boundary", "each agent's output is parsed against a Zod schema — malformed AI responses retry automatically (observed self-correcting live)"),
    ("Metered", "every AI call is token-accounted and billed to the wallet in credits, with per-user/org budget hard caps"),
], size=14.5, gap=9)
footer(s, 4)

# ═══ 5. Shorts Studio flow ═════════════════════════════════════════════════════
s = slide(); header(s, "Core flow · repurposing", "Shorts Studio: one long-form video becomes a content week")
stages2 = [
    ("Import", "yt-dlp source\n+ captions", ACCENT2), ("Analyze", "transcript, scenes,\ntopics, chapters", ACCENT2),
    ("Highlights", "AI-scored viral\nmoments", GOLD), ("Clip", "9:16 / 1:1 timeline\nper platform", ACCENT2),
    ("Track + Reframe", "face + motion\ntracked crop", ACCENT), ("Render", "captions burned,\nGPU-aware", ACCENT2),
    ("Approve + Publish", "human gate →\nYouTube", GREEN),
]
x = Inches(0.45); y = Inches(1.95)
for i, (name, sub, col) in enumerate(stages2):
    w = Inches(1.65)
    card(s, x, y, w, Inches(1.45), fill=PANEL, line_color=col)
    txt(s, x + Inches(0.05), y + Inches(0.12), Inches(1.55), Inches(0.55), name, size=13, color=col, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.05), y + Inches(0.66), Inches(1.55), Inches(0.7), sub, size=10, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < len(stages2) - 1:
        arrow(s, x + w + Inches(0.02), y + Inches(0.58), w=Inches(0.14))
    x += w + Inches(0.18)
c = card(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(2.85), line_color=ACCENT)
txt(s, Inches(0.85), Inches(4.05), Inches(5.5), Inches(0.4), "SUBJECT-TRACKED REFRAMING (new)", size=13, color=ACCENT, bold=True)
bullets(s, Inches(0.85), Inches(4.5), Inches(5.5), Inches(2.1), [
    "On-device face detection (pico cascade) finds the speaker in every sampled frame",
    "Movement-centroid fallback tracks robots, demos, screen action",
    "Verified live: 88/88 frames face-locked on a presenter video; crop pans with the subject",
], size=12.5, gap=7)
c = card(s, Inches(6.85), Inches(3.85), Inches(5.85), Inches(2.85))
txt(s, Inches(7.1), Inches(4.05), Inches(5.3), Inches(0.4), "BEYOND SHORTS", size=13, color=ACCENT2, bold=True)
bullets(s, Inches(7.1), Inches(4.5), Inches(5.4), Inches(2.1), [
    ("Social factory", "quote cards, carousels, blog posts, newsletters from the same transcript"),
    ("Chapter sync", "AI-detected chapters pushed to the YouTube description"),
    ("Semantic search", "ask the library a question, jump to the exact clip moment"),
], size=12.5, gap=7, lead_color=ACCENT2)
footer(s, 5)

# ═══ 6. Feature map ════════════════════════════════════════════════════════════
s = slide(); header(s, "Features & functions", "Feature map by module")
cols = [
    ("CREATE", ACCENT, [
        "Full long-form agent pipeline",
        "Shorts Studio (6 clip formats)",
        "Standalone multi-track video editor — filters, transitions, keyframes, text, audio mixing, export presets",
        "AI thumbnails + quote cards",
    ]),
    ("GROW", ACCENT2, [
        "SEO & metadata engine",
        "Channel analytics + BI module",
        "Trend & topic research agents",
        "AI copilot with editing commands",
        "Semantic library search",
    ]),
    ("OPERATE", GOLD, [
        "Per-channel automation with daily quotas",
        "Approvals inbox (human gate)",
        "Orgs & teams, shared wallets, budget caps",
        "Referrals, trials, marketplace offers",
    ]),
    ("EXTEND", GREEN, [
        "Developer portal: scoped API keys",
        "Webhooks with signed deliveries",
        "Public REST API (rate-limited per key)",
        "n8n workflow definitions",
    ]),
]
x = Inches(0.6)
for name, col, items in cols:
    card(s, x, Inches(1.75), Inches(2.98), Inches(4.9))
    chip(s, x + Inches(0.25), Inches(1.98), Inches(1.5), Inches(0.4), name, fill=col, tcolor=INK if col in (GOLD, GREEN, ACCENT2) else WHITE)
    bullets(s, x + Inches(0.22), Inches(2.55), Inches(2.62), Inches(3.9), items, size=11.5, gap=7, lead_color=col)
    x += Inches(3.14)
footer(s, 6)

# ═══ 7. Tools & AI stack ═══════════════════════════════════════════════════════
s = slide(); header(s, "Tools & intelligence", "The agent workforce and the tools it drives")
c = card(s, Inches(0.6), Inches(1.75), Inches(6.1), Inches(4.9))
txt(s, Inches(0.85), Inches(1.95), Inches(5.6), Inches(0.4), "AI AGENTS (per-role specialists)", size=13, color=ACCENT, bold=True)
bullets(s, Inches(0.85), Inches(2.4), Inches(5.6), Inches(4.1), [
    ("Supervisor", "orchestrates every pipeline; agents stay stateless + idempotent"),
    ("Research / FactCheck", "no fabricated facts — claims trace to captured sources"),
    ("Compliance", "9-category audit (copyright, misinformation, advertiser-friendliness…) with BLOCK-severity flags"),
    ("Script / SEO / Metadata", "channel voice profiles + brand kits"),
    ("Voice / Image / Music / Video / Subtitle", "media generation with provenance metadata"),
    ("EditPlan / Highlight / Chapter", "turns analysis into timelines and clips"),
], size=12.5, gap=8)
c = card(s, Inches(6.95), Inches(1.75), Inches(5.75), Inches(4.9))
txt(s, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.4), "TOOLING", size=13, color=ACCENT2, bold=True)
bullets(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(4.1), [
    ("Multi-provider AI client", "Claude (primary) + GPT + Gemini with health-scored failover, retries, response caching, cost accounting per call"),
    ("ffmpeg", "renders, reframes, caption burn-in — GPU (NVENC) aware with CPU fallback"),
    ("yt-dlp", "source acquisition + auto-captions"),
    ("pico face detection", "on-device, zero API cost, MIT-licensed cascade"),
    ("YouTube Data API", "library browse, metadata, chapters, uploads"),
    ("Stripe", "subscriptions + credit recharges, idempotent webhook settlement"),
], size=12.5, gap=8, lead_color=ACCENT2)
footer(s, 7)

# ═══ 8. Architecture & trust ═══════════════════════════════════════════════════
s = slide(); header(s, "Architecture & trust", "Production-grade engineering under the hood")
c = card(s, Inches(0.6), Inches(1.75), Inches(6.1), Inches(4.9))
txt(s, Inches(0.85), Inches(1.95), Inches(5.6), Inches(0.4), "PLATFORM", size=13, color=ACCENT2, bold=True)
bullets(s, Inches(0.85), Inches(2.4), Inches(5.6), Inches(4.1), [
    ("Stack", "TypeScript monorepo — NestJS API (40+ modules), Next.js web, Postgres, Redis, BullMQ job queue, socket.io live progress"),
    ("Scale-ready", "Kubernetes architecture shipped: API autoscales 2→8 pods, Redis-shared rate limits and caches, migrations per rollout"),
    ("CI/CD", "lint, typecheck, 650+ unit tests, SAST (Semgrep), DAST (OWASP ZAP), 3-browser E2E, dependency audit — every push"),
    ("Observability", "Prometheus metrics, Grafana dashboards, Sentry, structured logs with correlation IDs"),
], size=12.5, gap=9)
c = card(s, Inches(6.95), Inches(1.75), Inches(5.75), Inches(4.9), line_color=GREEN)
txt(s, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.4), "SECURITY & SAFETY (verified live)", size=13, color=GREEN, bold=True)
bullets(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(4.1), [
    ("Auth", "JWT + rotating refresh tokens — 10-point attack test passed (tampering, alg:none, replay, logout)"),
    ("Rate limiting", "Redis-backed on auth + API keys, multi-instance safe"),
    ("Secrets", "OAuth tokens AES-256-GCM encrypted at rest; boot refuses weak keys in production"),
    ("Billing integrity", "webhook settlement idempotent at event AND payment level — replay-tested, no double-credit"),
    ("Compliance gate", "cached verdicts re-checked on every enforce; cache can never soften the gate"),
], size=12.5, gap=9, lead_color=GREEN)
footer(s, 8)

# ═══ 9. Business model ═════════════════════════════════════════════════════════
s = slide(); header(s, "Business model", "Subscriptions for access, credits for consumption")
tiers = [
    ("FREE", "Trial credits on signup, full pipeline access, watermark-level quotas", MUTED),
    ("STARTER", "Solo creators: monthly credit quota, 1–2 channels, automation basics", ACCENT2),
    ("PRO", "Full-time creators: higher quotas, all formats, priority rendering", ACCENT),
    ("AGENCY", "Teams & orgs: shared wallets, budget hard-caps per client, seats + RBAC", GOLD),
]
x = Inches(0.6)
for name, desc, col in tiers:
    card(s, x, Inches(1.8), Inches(2.98), Inches(2.2), line_color=col)
    txt(s, x + Inches(0.22), Inches(2.0), Inches(2.5), Inches(0.45), name, size=19, color=col, bold=True)
    txt(s, x + Inches(0.22), Inches(2.5), Inches(2.55), Inches(1.4), desc, size=12, color=MUTED, line_spacing=1.15)
    x += Inches(3.14)
bullets(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.3), [
    ("Credits meter real cost", "every AI/video/music job debits the wallet (default 100 credits ≈ $1); margins hold because provider costs are accounted per call"),
    ("Recharge on demand", "Stripe checkout for credit top-ups and marketplace packs — settlement flow fully tested end to end, replay-safe"),
    ("Expansion built in", "referral credits, trial-to-paid conversion engine, first-recharge rewards, org budget periods"),
    ("Anti-abuse", "device-fingerprint trial scoring, recharge freeze on dispute, fraud-hold path"),
], size=14, gap=8)
footer(s, 9)

# ═══ 10. Target users ══════════════════════════════════════════════════════════
s = slide(); header(s, "Market strategy · who", "Clear target users, in order of attack")
segs = [
    ("1 · The Serious Solo Creator", "10K–500K subs, monetized, 1–2 uploads/week",
     "Pain: production time caps output; repurposing never happens.\nBuys: PRO. Hook: Shorts Studio + automation.", ACCENT),
    ("2 · The Expert-on-Camera", "coaches, educators, developers, finance/legal pros",
     "Pain: knows the subject, hates the toolchain; terrified of compliance strikes.\nBuys: STARTER→PRO. Hook: fact-checked scripts + face-tracked shorts.", ACCENT2),
    ("3 · The Creator Agency", "manages 5–50 channels for clients",
     "Pain: per-client budgets, approvals, and consistency at scale.\nBuys: AGENCY. Hook: orgs, shared wallets, budget hard-caps, human approval workflow.", GOLD),
]
x = Inches(0.6)
for title, who, body, col in segs:
    card(s, x, Inches(1.8), Inches(4.0), Inches(4.4), line_color=col)
    txt(s, x + Inches(0.25), Inches(2.0), Inches(3.5), Inches(0.8), title, size=16.5, color=col, bold=True, line_spacing=1.0)
    txt(s, x + Inches(0.25), Inches(2.75), Inches(3.5), Inches(0.6), who, size=12, color=WHITE, line_spacing=1.05)
    txt(s, x + Inches(0.25), Inches(3.45), Inches(3.5), Inches(2.6), body, size=12, color=MUTED, line_spacing=1.2)
    x += Inches(4.18)
txt(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.5),
    "Explicitly NOT targeting: content farms, faceless-spam channels, engagement-bait operations — excluded by product design, not just positioning.",
    size=12.5, color=MUTED)
footer(s, 10)

# ═══ 11. Value proposition ═════════════════════════════════════════════════════
s = slide(); header(s, "Market strategy · why", "Value proposition")
c = card(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.5), fill=PANEL, line_color=ACCENT)
tf = txt(s, Inches(0.95), Inches(2.1), Inches(11.4), Inches(1.0),
    "Publish 5× more content from the same ideas — without risking the channel that feeds you.",
    size=23, color=WHITE, bold=True, line_spacing=1.05)
rows = [
    ("For the creator", "One video's effort becomes a week of content: the long-form, 3–5 tracked shorts, quote cards, a blog post, a newsletter — all from one pipeline."),
    ("For the channel", "Monetization-safe by construction: fact-checked claims, a 9-category compliance audit that can block, and a human finger on every publish button."),
    ("For the business", "Costs are metered, budgeted, and hard-capped in credits — no surprise AI bills; agencies allocate per client."),
]
y = Inches(3.7)
for lead, body in rows:
    txt(s, Inches(0.7), y, Inches(2.6), Inches(0.5), lead, size=15, color=ACCENT2, bold=True)
    txt(s, Inches(3.4), y, Inches(9.3), Inches(0.9), body, size=14, color=MUTED, line_spacing=1.15)
    y += Inches(0.98)
footer(s, 11)

# ═══ 12. Positioning ═══════════════════════════════════════════════════════════
s = slide(); header(s, "Market strategy · where", "Positioning: own the full workflow + trust")
tbl_data = [
    ("", "Opus Clip / 2short", "vidIQ / TubeBuddy", "Descript / CapCut", "CreatorForce"),
    ("Clips & repurposing", "●", "—", "◐", "●"),
    ("SEO & channel growth", "—", "●", "—", "●"),
    ("Full editing", "—", "—", "●", "●"),
    ("Idea→publish pipeline", "—", "—", "—", "●"),
    ("Compliance / fact-check gate", "—", "—", "—", "●"),
    ("Team budgets & approvals", "—", "—", "◐", "●"),
]
rows_n, cols_n = len(tbl_data), 5
tbl = s.shapes.add_table(rows_n, cols_n, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.3)).table
tbl.columns[0].width = Inches(3.4)
for ci in range(1, 5):
    tbl.columns[ci].width = Inches(2.175)
for ri, row in enumerate(tbl_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if ri else INK
        if ci == 4 and ri:
            cell.fill.fore_color.rgb = RGBColor(0x30, 0x1B, 0x24)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(13); r.font.name = "Segoe UI"
        r.font.bold = (ri == 0 or ci == 0)
        r.font.color.rgb = WHITE if ri == 0 else (GREEN if val == "●" and ci == 4 else (ACCENT2 if val == "●" else MUTED))
txt(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.4),
    "Positioning statement:  For serious YouTube creators and the agencies behind them, CreatorForce is the "
    "content operating system that turns one idea into a week of published, monetization-safe content — unlike "
    "point tools that optimize a single step and leave the risk and the stitching to you.",
    size=15.5, color=WHITE, line_spacing=1.25)
footer(s, 12)

# ═══ 13. Messaging ═════════════════════════════════════════════════════════════
s = slide(); header(s, "Market strategy · what we say", "Messaging house")
c = card(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(1.0), line_color=GOLD)
txt(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.6),
    "“Your channel, on autopilot. Your judgment, in charge.”", size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
pillars = [
    ("SPEED", "“One video in. A content week out.”",
     ["Pipeline demo: idea → rendered video", "Shorts Studio before/after timelines", "Automation quotas doing the boring part"], ACCENT),
    ("SAFETY", "“AI that won't get you demonetized.”",
     ["Compliance audit screenshots w/ BLOCK flags", "Fact-check with source tracing", "Human approval inbox — nothing ships itself"], GREEN),
    ("CRAFT", "“Shorts that follow the story — literally.”",
     ["Face-tracked reframing demos", "Multi-track editor + brand kits", "Channel-voice scripts, not generic AI slop"], ACCENT2),
]
x = Inches(0.6)
for name, tagline, proofs, col in pillars:
    card(s, x, Inches(3.05), Inches(4.0), Inches(3.5))
    chip(s, x + Inches(0.25), Inches(3.25), Inches(1.3), Inches(0.4), name, fill=col, tcolor=INK if col in (GREEN, ACCENT2) else WHITE)
    txt(s, x + Inches(0.25), Inches(3.8), Inches(3.55), Inches(0.75), tagline, size=14.5, color=WHITE, bold=True, line_spacing=1.05)
    bullets(s, x + Inches(0.25), Inches(4.6), Inches(3.55), Inches(1.8), proofs, size=11.5, gap=6, lead_color=col)
    x += Inches(4.18)
footer(s, 13)

# ═══ 14. GTM plan ══════════════════════════════════════════════════════════════
s = slide(); header(s, "Market strategy · how", "Go-to-market: three phases, each funds the next")
phases = [
    ("PHASE 1 · Prove (0–3 mo)", ACCENT, [
        "Closed beta: 50–100 hand-picked creators in 2 niches (dev/education, faith & community — church pack already built)",
        "White-glove onboarding; weekly ship cadence from feedback",
        "Success gate: 30% weekly active, 5 documented “content week” case studies, <5% compliance-gate false blocks",
    ]),
    ("PHASE 2 · Product-led growth (3–9 mo)", ACCENT2, [
        "Self-serve FREE→STARTER with trial credits; first-recharge reward live",
        "Creators ARE the channel: case-study videos, before/after shorts, referral credits (built in)",
        "Content moat: SEO plays on “YouTube compliance”, “shorts face tracking”, “AI video pipeline”",
        "Success gate: CAC < 1 month of PRO; trial→paid ≥ 8%",
    ]),
    ("PHASE 3 · Agencies & platform (9–18 mo)", GOLD, [
        "AGENCY tier outbound: agencies discovered via multi-channel usage patterns",
        "Developer API GA + marketplace packs; partner white-labeling (roadmap Phase 7)",
        "Success gate: 20% of MRR from orgs; net revenue retention > 110%",
    ]),
]
y = Inches(1.7)
for name, col, items in phases:
    card(s, Inches(0.6), y, Inches(12.1), Inches(1.5), line_color=col)
    txt(s, Inches(0.85), y + Inches(0.08), Inches(3.6), Inches(0.5), name, size=15, color=col, bold=True)
    bullets(s, Inches(4.3), y + Inches(0.06), Inches(8.2), Inches(1.4), items, size=10.5, gap=2, lead_color=col)
    y += Inches(1.62)
txt(s, Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.4),
    "Credibility check: every growth mechanism named above (referrals, trials, rewards, orgs, API keys) is already built and tested — GTM is activation, not construction.",
    size=12, color=GREEN)
footer(s, 14)

# ═══ 15. Roadmap ═══════════════════════════════════════════════════════════════
s = slide(); header(s, "Where it goes next", "Roadmap: phases 4–7")
roadmap = [
    ("NOW · Phase 4–5 finish", ACCENT, "Cloud storage (R2) · external video/music providers (Veo, Runway, Suno…) · render→publish hand-off · staging env · k6/ZAP go-live runbook"),
    ("NEXT · Phase 6", ACCENT2, "AI autonomy: copilot orchestrating full workflows · cross-channel intelligence · auto content calendars · real-time trends"),
    ("LATER · Phase 7", GOLD, "Platform: agency white-labeling · more AI providers (DeepSeek, Grok, Mistral…) · i18n · public developer API GA"),
]
y = Inches(2.0)
for name, col, body in roadmap:
    chip(s, Inches(0.6), y, Inches(2.6), Inches(0.5), name, fill=col, tcolor=INK if col in (ACCENT2, GOLD) else WHITE)
    txt(s, Inches(3.5), y + Inches(0.02), Inches(9.2), Inches(1.1), body, size=14, color=MUTED, line_spacing=1.2)
    y += Inches(1.45)
c = card(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.75), line_color=GREEN)
txt(s, Inches(0.85), Inches(6.26), Inches(11.6), Inches(0.5),
    "Status today: core platform, Shorts Studio, editor, automation, billing and K8s architecture are built, tested (652 unit tests, CI green) and live-verified end to end.",
    size=13.5, color=GREEN)
footer(s, 15)

# ═══ 16. Close ═════════════════════════════════════════════════════════════════
s = slide()
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.5), Inches(3.5), Inches(8), Inches(8))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x25, 0x14, 0x1E); glow.line.fill.background(); glow.shadow.inherit = False
txt(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.2),
    "AI CreatorForce", size=44, bold=True)
txt(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(0.8),
    "The workflow is the product. The trust is the moat.", size=22, color=ACCENT)
bullets(s, Inches(0.7), Inches(4.3), Inches(11.5), Inches(2.0), [
    "Software: a full YouTube content OS — agents, studio, editor, automation, API",
    "Market: serious creators and agencies who can't afford to get it wrong",
    "Plan: prove with 100 creators → grow product-led → land agencies on the tier built for them",
], size=16, gap=10)
txt(s, Inches(0.7), Inches(6.7), Inches(11), Inches(0.4),
    "aicreatorforce.net  ·  July 2026", size=13, color=MUTED)

prs.save("docs/CreatorForce-Product-Overview.pptx")
print("Saved docs/CreatorForce-Product-Overview.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
